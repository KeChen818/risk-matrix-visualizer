import numpy as np
import matplotlib.pyplot as plt
import textwrap
from pptx import Presentation
from pptx.util import Inches

def plot_risk_matrix_top_left(risks, material_grids, likelihood_labels, impact_labels, save_path="risk_matrix.png"):
    """
    Plots a risk matrix with bubbles indexed by risk IDs and 
    configurable material (highlighted) grids, supporting 6x6 (36 risks per cell).
    Saves the plot as an image for use in slides.
    """
    grid_size = 6
    cols = 6  # Maximum columns per row in each grid cell
    rows = 6  # Maximum rows per grid cell
    bubble_radius = 0.05  # Bubble size
    bubble_spacing = 0.12  # Bubble spacing for correct layout

    # Wrap labels for better readability
    wrapped_likelihood_labels = [textwrap.fill(label, width=15) for label in likelihood_labels]
    wrapped_impact_labels = [textwrap.fill(label, width=15) for label in impact_labels]

    # Organize risks by grid cell
    risk_counts = {}
    risks_in_cell = {}
    for risk_id, x, y in risks:
        key = (x - 1, y - 1)  # Convert to 0-based index
        risk_counts[key] = risk_counts.get(key, 0) + 1
        risks_in_cell.setdefault(key, []).append(risk_id)  # Store risk IDs per cell

    # Create plot
    fig, ax = plt.subplots(figsize=(8, 8))
    ax.set_xlim(-0.5, grid_size - 0.5)
    ax.set_ylim(-0.5, grid_size - 0.5)
    ax.set_xticks(np.arange(grid_size))
    ax.set_yticks(np.arange(grid_size))
    ax.grid(True, linestyle="--", linewidth=0.5)
    ax.set_aspect('equal')  # Keep square aspect

    # Highlight material grids
    for (x, y) in material_grids:
        ax.add_patch(plt.Rectangle((x - 0.5, y - 0.5), 1, 1, color="lightgrey", alpha=0.5))  # Light grey background
        ax.text(x, y, "Material", fontsize=10, ha="center", va="center", color="white", weight="bold")

    # Plot bubbles in each grid cell with assigned risk IDs (6x6 max per cell)
    for (x, y), risks in risks_in_cell.items():
        col_start = x - 0.4  # Ensure bubbles start exactly from the top-left corner of the grid cell
        row_start = y - 0.4  # Adjust row_start to align bubbles to the top

        # Arrange bubbles: fill each row first before moving downward, staying in the same grid cell
        for i, risk_id in enumerate(risks[:rows * cols]):  # Limit to 36 per cell
            row_offset = (i // cols) * bubble_spacing  # Moves downward correctly now
            col_offset = (i % cols) * bubble_spacing  # Fill **left to right** in a row
            px, py = col_start + col_offset, row_start + row_offset

            # Plot bubbles
            ax.add_patch(plt.Circle((px, py), bubble_radius, color='r'))

            # Add risk ID text inside the bubble
            ax.text(px, py, str(risk_id), fontsize=4, ha='center', va='center', color='white', weight='bold')

    # Invert y-axis for correct risk matrix orientation
    plt.gca().invert_yaxis()
    plt.xlabel("Likelihood", fontsize=12, weight="bold", labelpad=15)
    plt.ylabel("Impact", fontsize=12, weight="bold", labelpad=15)
    plt.xticks(np.arange(0.5, grid_size, 1), wrapped_likelihood_labels, rotation=0, ha="center", fontsize=9)
    plt.yticks(np.arange(0.5, grid_size, 1), reversed(wrapped_impact_labels), fontsize=9)
    plt.title("Risk Matrix with Wrapped Labels and Top-Left Aligned Bubbles")
    
    # Save plot as image
    plt.savefig(save_path, bbox_inches='tight')
    plt.close()

    return save_path


def create_ppt_with_risk_matrix(image_path, ppt_path="risk_matrix_presentation.pptx"):
    """
    Creates a PowerPoint slide with the risk matrix image taking up the right 2/5 portion of the slide.
    """
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Define image position (Right 2/5 of the slide)
    img_left = Inches(4)  # Adjusted for 2/5 right alignment
    img_top = Inches(1)
    img_width = Inches(6)  # Scale image size properly
    img_height = Inches(5)
    
    # Add image to slide
    slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)
    
    # Save PowerPoint
    prs.save(ppt_path)
    return ppt_path

# Example Fake Data (100 Risks with IDs, Likelihood, Impact)
num_risks = 100
risks = [(i+1, np.random.randint(1, 7), np.random.randint(1, 7)) for i in range(num_risks)]

# Define Material Grid Positions (Example: First two rows)
material_grids = [(x, y) for x in range(6) for y in range(2)]  # Top two rows

# Likelihood and Impact Labels
likelihood_labels = [
    "1-Low (once in 20 years)",
    "2-Minor (once in 10 years)",
    "3-Moderate (once in 5 years)",
    "4-Significant (once in 2 years)",
    "5-Major (once a year)",
    "6-Critical (multiple times a year)"
]

impact_labels = [
    "1-Low (<$1M)",
    "2-Minor ($1M-$10M)",
    "3-Moderate ($10M-$50M)",
    "4-Significant ($50M-$500M)",
    "5-High (>$500M)",
    "6-Severe (>$1B)"
]

# Generate risk matrix image
image_path = plot_risk_matrix_top_left(risks, material_grids, likelihood_labels, impact_labels)

# Create PowerPoint with risk matrix in right 2/5 portion
ppt_path = create_ppt_with_risk_matrix(image_path)
print(f"PowerPoint saved at: {ppt_path}")

